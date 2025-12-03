"""
Generate Awesome Presentation for Hyperparameter Sensitivity Analysis
CSIT332 - Principles of Machine Learning
"""

import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Set style
sns.set_style('whitegrid')
plt.rcParams['figure.dpi'] = 150

# Load all sensitivity results
datasets = {
    'Census Income': 'outputs/results/census_income_sensitivity.csv',
    'Chronic Kidney Disease': 'outputs/results/ckd_sensitivity.csv',
    'Energy Efficiency': 'outputs/results/energy_sensitivity.csv',
    'Thyroid Disease': 'outputs/results/thyroid_sensitivity.csv',
    'Wholesale Customers': 'outputs/results/wholesale_sensitivity.csv'
}

sensitivity_data = {}
for name, path in datasets.items():
    sensitivity_data[name] = pd.read_csv(path)

print("✓ Loaded all sensitivity data")

# Create presentation
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

# Define colors
BLUE = RGBColor(0, 112, 192)
ORANGE = RGBColor(255, 127, 0)
GREEN = RGBColor(0, 176, 80)
RED = RGBColor(192, 0, 0)
DARK_GRAY = RGBColor(64, 64, 64)

def add_title_slide(prs, title, subtitle):
    """Add title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Background color
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(240, 248, 255)
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = BLUE
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(3.8), Inches(8), Inches(1.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = DARK_GRAY
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    return slide

def add_section_slide(prs, title):
    """Add section divider slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BLUE
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(1.5))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, content_type='bullet'):
    """Add content slide with title"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = BLUE
    
    return slide

# Slide 1: Title
add_title_slide(prs, 
    "Hyperparameter Sensitivity Trends Across Dataset Types",
    "UCI Machine Learning Repository Analysis\nCSIT332 - Principles of Machine Learning")

# Slide 2: Project Overview
slide = add_content_slide(prs, "Project Overview")
content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
tf = content_box.text_frame
tf.word_wrap = True

bullets = [
    ("Objective", "Analyze how different classifiers respond to hyperparameter changes across diverse datasets"),
    ("Datasets Analyzed", "5 UCI ML datasets: Census Income, Chronic Kidney Disease, Energy Efficiency, Thyroid Disease, Wholesale Customers"),
    ("Classifiers Tested", "6 algorithms: KNN, Logistic Regression, SVM, MLP, Decision Tree, Naive Bayes"),
    ("Methodology", "GridSearchCV with 3-fold cross-validation to measure variance in accuracy across hyperparameter combinations"),
    ("Key Metric", "Variance of accuracy scores - higher variance = more sensitive to hyperparameter tuning")
]

for label, text in bullets:
    p = tf.add_paragraph()
    p.text = f"• {label}: "
    p.font.size = Pt(18)
    p.font.color.rgb = DARK_GRAY
    p.level = 0
    
    run = p.runs[0]
    run.font.bold = True
    run.font.color.rgb = BLUE
    
    p.text += text
    p.space_after = Pt(12)

# Slide 3: Methodology
slide = add_content_slide(prs, "Methodology: Measuring Hyperparameter Sensitivity")
content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
tf = content_box.text_frame

process = [
    "1. Define comprehensive hyperparameter grids for each classifier",
    "2. GridSearchCV tests all combinations (28 for KNN, 48 for Logistic Regression, etc.)",
    "3. Record cross-validation accuracy for each combination",
    "4. Calculate variance, standard deviation, and range of accuracy scores",
    "5. Rank classifiers by sensitivity (variance) for each dataset"
]

for item in process:
    p = tf.add_paragraph()
    p.text = item
    p.font.size = Pt(22)
    p.font.color.rgb = DARK_GRAY
    p.space_after = Pt(18)

formula_box = slide.shapes.add_textbox(Inches(2), Inches(5.5), Inches(6), Inches(1))
formula_tf = formula_box.text_frame
p = formula_tf.paragraphs[0]
p.text = "Variance = measure of spread in accuracy scores"
p.font.size = Pt(20)
p.font.italic = True
p.font.color.rgb = ORANGE
p.alignment = PP_ALIGN.CENTER

# Slide 4: Section - Results
add_section_slide(prs, "Key Findings")

# Slide 5: Overall Sensitivity Heatmap
slide = add_content_slide(prs, "Sensitivity Heatmap Across All Datasets")

# Create heatmap
classifiers = ['KNN', 'Logistic Regression', 'SVM', 'MLP', 'Decision Tree', 'Naive Bayes']
dataset_names = list(sensitivity_data.keys())

# Build matrix
heatmap_data = []
for clf in classifiers:
    row = []
    for dataset in dataset_names:
        df = sensitivity_data[dataset]
        variance = df[df['Classifier'] == clf]['Variance'].values[0] if clf in df['Classifier'].values else 0
        row.append(variance)
    heatmap_data.append(row)

heatmap_df = pd.DataFrame(heatmap_data, index=classifiers, columns=['Census', 'CKD', 'Energy', 'Thyroid', 'Wholesale'])

# Plot
fig, ax = plt.subplots(figsize=(10, 6))
sns.heatmap(heatmap_df, annot=True, fmt='.4f', cmap='YlOrRd', cbar_kws={'label': 'Variance'}, 
            linewidths=0.5, ax=ax)
ax.set_title('Hyperparameter Sensitivity by Classifier and Dataset', fontsize=16, fontweight='bold', pad=20)
ax.set_xlabel('Dataset', fontsize=12, fontweight='bold')
ax.set_ylabel('Classifier', fontsize=12, fontweight='bold')
plt.tight_layout()
plt.savefig('temp_heatmap.png', dpi=150, bbox_inches='tight')
plt.close()

slide.shapes.add_picture('temp_heatmap.png', Inches(0.8), Inches(1.8), width=Inches(8.4))
os.remove('temp_heatmap.png')

# Slide 6: Most Sensitive Classifiers
slide = add_content_slide(prs, "Most Sensitive Classifiers (Require Careful Tuning)")

# Get top sensitive from each dataset
top_sensitive = []
for dataset, df in sensitivity_data.items():
    top = df.nlargest(1, 'Variance')
    top_sensitive.append({
        'dataset': dataset,
        'classifier': top['Classifier'].values[0],
        'variance': top['Variance'].values[0],
        'range': top['Range'].values[0]
    })

top_sensitive_df = pd.DataFrame(top_sensitive).sort_values('variance', ascending=False)

# Create bar chart
fig, ax = plt.subplots(figsize=(10, 6))
colors = ['#FF6B6B', '#FF8C42', '#FFA62B', '#FFD93D', '#6BCB77']
bars = ax.barh(range(len(top_sensitive_df)), top_sensitive_df['variance'], color=colors)
ax.set_yticks(range(len(top_sensitive_df)))
ax.set_yticklabels([f"{row['dataset']}\n({row['classifier']})" 
                     for _, row in top_sensitive_df.iterrows()], fontsize=11)
ax.set_xlabel('Variance (Sensitivity)', fontsize=12, fontweight='bold')
ax.set_title('Most Sensitive Classifier per Dataset', fontsize=14, fontweight='bold', pad=15)
ax.grid(axis='x', alpha=0.3)

# Add value labels
for i, (idx, row) in enumerate(top_sensitive_df.iterrows()):
    ax.text(row['variance'] + 0.001, i, f"{row['variance']:.4f}\n(Range: {row['range']:.2%})", 
            va='center', fontsize=10, fontweight='bold')

plt.tight_layout()
plt.savefig('temp_sensitive.png', dpi=150, bbox_inches='tight')
plt.close()

slide.shapes.add_picture('temp_sensitive.png', Inches(0.8), Inches(1.8), width=Inches(8.4))
os.remove('temp_sensitive.png')

# Slide 7: Most Robust Classifiers
slide = add_content_slide(prs, "Most Robust Classifiers (Work Well with Defaults)")

# Get least sensitive from each dataset
least_sensitive = []
for dataset, df in sensitivity_data.items():
    bottom = df.nsmallest(1, 'Variance')
    least_sensitive.append({
        'dataset': dataset,
        'classifier': bottom['Classifier'].values[0],
        'variance': bottom['Variance'].values[0],
        'accuracy': bottom['Test Accuracy'].values[0]
    })

least_sensitive_df = pd.DataFrame(least_sensitive).sort_values('variance')

# Create chart
fig, ax = plt.subplots(figsize=(10, 6))
colors = ['#4ECDC4', '#44A08D', '#6C5B7B', '#355C7D', '#2A9D8F']
bars = ax.barh(range(len(least_sensitive_df)), least_sensitive_df['variance'], color=colors)
ax.set_yticks(range(len(least_sensitive_df)))
ax.set_yticklabels([f"{row['dataset']}\n({row['classifier']})" 
                     for _, row in least_sensitive_df.iterrows()], fontsize=11)
ax.set_xlabel('Variance (Lower = More Robust)', fontsize=12, fontweight='bold')
ax.set_title('Most Robust Classifier per Dataset', fontsize=14, fontweight='bold', pad=15)
ax.grid(axis='x', alpha=0.3)

# Add labels
for i, (idx, row) in enumerate(least_sensitive_df.iterrows()):
    variance_text = f"{row['variance']:.6f}" if row['variance'] > 0 else "0.0000"
    ax.text(row['variance'] + 0.00005, i, f"{variance_text}\n(Acc: {row['accuracy']:.2%})", 
            va='center', fontsize=10, fontweight='bold')

plt.tight_layout()
plt.savefig('temp_robust.png', dpi=150, bbox_inches='tight')
plt.close()

slide.shapes.add_picture('temp_robust.png', Inches(0.8), Inches(1.8), width=Inches(8.4))
os.remove('temp_robust.png')

# Slide 8: Dataset Characteristics
slide = add_content_slide(prs, "Key Insights by Dataset")

# Compile insights
insights = []
for dataset, df in sensitivity_data.items():
    most_sensitive = df.nlargest(1, 'Variance')
    least_sensitive = df.nsmallest(1, 'Variance')
    avg_variance = df['Variance'].mean()
    
    insights.append({
        'dataset': dataset,
        'avg_variance': avg_variance,
        'most': most_sensitive['Classifier'].values[0],
        'least': least_sensitive['Classifier'].values[0],
        'best_acc': df['Test Accuracy'].max()
    })

insights_df = pd.DataFrame(insights).sort_values('avg_variance', ascending=False)

# Create comparison chart
fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(12, 5))

# Average variance by dataset
ax1.bar(range(len(insights_df)), insights_df['avg_variance'], 
        color=['#E63946', '#F77F00', '#FCBF49', '#06A77D', '#118AB2'])
ax1.set_xticks(range(len(insights_df)))
ax1.set_xticklabels([d.split()[0] for d in insights_df['dataset']], rotation=45, ha='right')
ax1.set_ylabel('Average Variance', fontweight='bold')
ax1.set_title('Dataset Difficulty\n(Higher = More Sensitive Overall)', fontweight='bold', pad=10)
ax1.grid(axis='y', alpha=0.3)

# Best accuracy by dataset  
ax2.bar(range(len(insights_df)), insights_df['best_acc'], 
        color=['#06A77D', '#118AB2', '#073B4C', '#FCBF49', '#F77F00'])
ax2.set_xticks(range(len(insights_df)))
ax2.set_xticklabels([d.split()[0] for d in insights_df['dataset']], rotation=45, ha='right')
ax2.set_ylabel('Best Test Accuracy', fontweight='bold')
ax2.set_title('Dataset Performance\n(Achievable Accuracy)', fontweight='bold', pad=10)
ax2.set_ylim([0.75, 1.0])
ax2.grid(axis='y', alpha=0.3)

plt.tight_layout()
plt.savefig('temp_insights.png', dpi=150, bbox_inches='tight')
plt.close()

slide.shapes.add_picture('temp_insights.png', Inches(0.5), Inches(2), width=Inches(9))
os.remove('temp_insights.png')

# Slide 9: Practical Recommendations
slide = add_content_slide(prs, "Practical Recommendations")
content_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.5), Inches(8.4), Inches(5.5))
tf = content_box.text_frame

recommendations = [
    ("For Medical Datasets (CKD, Thyroid)", "Logistic Regression & SVM are HIGHLY sensitive - invest time in hyperparameter tuning. MLP and Naive Bayes are robust alternatives."),
    ("For Demographic Data (Census)", "SVM shows highest sensitivity. All classifiers perform reasonably well with defaults except for SVM."),
    ("For Physical/Engineering Data (Energy)", "Logistic Regression and SVM need careful tuning. Tree-based methods are more forgiving."),
    ("For Business Data (Wholesale)", "SVM and Logistic Regression again show highest sensitivity. Consider MLP for robust performance."),
    ("General Rule", "Naive Bayes consistently shows near-zero variance - excellent starting point for baseline models!")
]

for title, text in recommendations:
    p = tf.add_paragraph()
    p.text = f"• {title}:\n  "
    p.font.size = Pt(16)
    p.level = 0
    
    run = p.runs[0]
    run.font.bold = True
    run.font.color.rgb = BLUE
    
    p.text += text
    p.font.color.rgb = DARK_GRAY
    p.space_after = Pt(14)

# Slide 10: Conclusion
slide = add_content_slide(prs, "Conclusions")
content_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(5))
tf = content_box.text_frame

conclusions = [
    "✓ Logistic Regression and SVM consistently show highest hyperparameter sensitivity across datasets",
    "✓ Naive Bayes is remarkably robust - almost zero variance in all cases",
    "✓ MLP shows very low sensitivity despite complex architecture",
    "✓ Medical datasets (CKD, Thyroid) require more careful tuning than other domains",
    "✓ Hyperparameter sensitivity varies significantly by dataset type",
    "✓ Understanding sensitivity helps practitioners allocate tuning effort effectively"
]

for conclusion in conclusions:
    p = tf.add_paragraph()
    p.text = conclusion
    p.font.size = Pt(20)
    p.font.color.rgb = DARK_GRAY
    p.space_after = Pt(16)
    if conclusion.startswith("✓"):
        p.runs[0].font.color.rgb = GREEN

# Add highlight box
highlight_box = slide.shapes.add_textbox(Inches(1.5), Inches(5.8), Inches(7), Inches(0.8))
highlight_tf = highlight_box.text_frame
p = highlight_tf.paragraphs[0]
p.text = "Key Takeaway: Dataset characteristics significantly impact hyperparameter sensitivity!"
p.font.size = Pt(18)
p.font.bold = True
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

# Add colored background to highlight box
from pptx.enum.shapes import MSO_SHAPE
shape = slide.shapes.add_shape(
    MSO_SHAPE.ROUNDED_RECTANGLE,
    Inches(1.5), Inches(5.8), Inches(7), Inches(0.8)
)
shape.fill.solid()
shape.fill.fore_color.rgb = ORANGE
shape.line.color.rgb = ORANGE
# Move text to front
slide.shapes._spTree.remove(highlight_box._element)
slide.shapes._spTree.insert(-1, highlight_box._element)

# Slide 11: Thank You
slide = add_title_slide(prs,
    "Thank You!",
    "Questions?\n\nHyperparameter Sensitivity Analysis Complete")

# Save presentation
output_path = 'outputs/Hyperparameter_Sensitivity_Analysis_Presentation.pptx'
prs.save(output_path)

print(f"\n{'='*80}")
print(f"✨ AWESOME PRESENTATION CREATED! ✨")
print(f"{'='*80}")
print(f"\nSaved to: {output_path}")
print(f"\nSlides created: {len(prs.slides)}")
print(f"\n✓ Title Slide")
print(f"✓ Project Overview")
print(f"✓ Methodology")
print(f"✓ Comprehensive Heatmap")
print(f"✓ Most Sensitive Classifiers")
print(f"✓ Most Robust Classifiers")  
print(f"✓ Dataset Insights")
print(f"✓ Practical Recommendations")
print(f"✓ Conclusions")
print(f"✓ Thank You Slide")
print(f"\n{'='*80}")
print("Ready to present! 🎉")
print(f"{'='*80}\n")
